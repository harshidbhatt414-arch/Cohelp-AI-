import streamlit as st
import pandas as pd
from pptx import Presentation
import io
from PIL import Image
import img2pdf
import time

# 1. Page Config
st.set_page_config(page_title="Cohelp AI Pro", layout="wide", page_icon="🤝")

# 2. State Management (This helps the buttons switch pages)
if 'page' not in st.session_state:
    st.session_state.page = "🏠 Home"

def set_page(page_name):
    st.session_state.page = page_name

# 3. Sidebar Navigation (Still available for quick switching)
with st.sidebar:
    st.image("https://img.icons8.com/clouds/200/handshake.png")
    st.title("Cohelp AI")
    if st.button("⬅️ Back to Home Screen"):
        set_page("🏠 Home")
    st.divider()
    st.caption("Founder: Harshid Bhatt")

# --- 🏠 MOBILE-STYLE HOME PAGE ---
if st.session_state.page == "🏠 Home":
    st.title("🤝 Cohelp AI Home")
    st.subheader("Select an app to get started")
    st.write("---")

    # Create the Grid (Like an App Drawer)
    col1, col2 = st.columns(2)
    col3, col4 = st.columns(2)

    with col1:
        st.image("https://img.icons8.com/fluency/96/microsoft-powerpoint.png", width=80)
        if st.button("Prompt to PPT", use_container_width=True):
            set_page("🪄 Prompt to PPT")

    with col2:
        st.image("https://img.icons8.com/fluency/96/gmail-new.png", width=80)
        if st.button("AI Email Pro", use_container_width=True):
            set_page("✉️ AI Email Pro")

    with col3:
        st.image("https://img.icons8.com/fluency/96/pdf-2.png", width=80)
        if st.button("Photo to PDF", use_container_width=True):
            set_page("📄 Photo to PDF")

    with col4:
        st.image("https://img.icons8.com/fluency/96/bar-chart.png", width=80)
        if st.button("Data Analysis", use_container_width=True):
            set_page("📊 Data Analysis")

    st.write("---")
    st.info("💡 Tip: You can always use the sidebar to jump between features.")

# --- 🪄 PROMPT TO PPT ---
elif st.session_state.page == "🪄 Prompt to PPT":
    st.header("🪄 Prompt to PPT")
    prompt = st.text_input("Enter your topic")
    photos = st.file_uploader("Add Photos", type=['jpg', 'png'], accept_multiple_files=True)
    if st.button("Generate PPT"):
        # (Same PPT Logic as before)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = prompt if prompt else "Business Plan"
        out = io.BytesIO()
        prs.save(out)
        st.download_button("Download PPT", out.getvalue(), "Cohelp.pptx")
    if st.button("Back"): set_page("🏠 Home")

# --- ✉️ AI EMAIL PRO ---
elif st.session_state.page == "✉️ AI Email Pro":
    st.header("✉️ Professional Email")
    reason = st.text_area("What is the reason?")
    if st.button("Generate"):
        st.code(f"To: Client\nSubject: Discussion\n\nDear Sir,\nRegarding {reason}...")
    if st.button("Back"): set_page("🏠 Home")

# --- 📄 PHOTO TO PDF ---
elif st.session_state.page == "📄 Photo to PDF":
    st.header("📄 Photo to PDF Scanner")
    scans = st.file_uploader("Upload photos", type=['jpg', 'png'], accept_multiple_files=True)
    if st.button("Create PDF"):
        if scans:
            pdf = img2pdf.convert([s.getvalue() for s in scans])
            st.download_button("Download PDF", pdf, "Scan.pdf")
    if st.button("Back"): set_page("🏠 Home")

# --- 📊 DATA ANALYSIS ---
elif st.session_state.page == "📊 Data Analysis":
    st.header("📊 Data Analysis")
    file = st.file_uploader("Upload CSV/Excel", type=['csv', 'xlsx'])
    if file:
        df = pd.read_csv(file) if file.name.endswith('csv') else pd.read_excel(file)
        st.dataframe(df)
    if st.button("Back"): set_page("🏠 Home")
      
